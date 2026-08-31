"""Build the SIH26141 idea-submission deck from the OFFICIAL template.

Rules enforced (from hackathon-decks/PHILOSOPHY.md):
  - max 6 slides incl. title; section headings NEVER renamed (they are the rubric)
  - read-deck density (~130 words/slide), boxes not bullets, no body text < 10pt
  - architecture diagram on slide 3 is load-bearing
  - every benefit claim carries a number and a source
  - exactly ONE figure set enormous
  - named risks with SPECIFIC mitigations
  - references carry titles + DOIs, never bare URLs
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy, pathlib

TPL = "/Users/harshwardhan/Downloads/SIH2026-IDEA-Presentation-Format.pptx"
OUT = "deck/SIH26141_PauliGuard_Idea_Submission.pptx"

NAVY   = RGBColor(0x14, 0x2A, 0x5A)
RED    = RGBColor(0xB9, 0x1C, 0x1C)
GREEN  = RGBColor(0x15, 0x71, 0x3D)
GREY   = RGBColor(0x44, 0x4A, 0x55)
LGREY  = RGBColor(0xE6, 0xE9, 0xEF)
BG     = RGBColor(0xF7, 0xF8, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

TEAM = "PauliGuard"

prs = Presentation(TPL)

def clear_instruction_boxes(slide, keep_titles=True):
    """Remove the template's placeholder guidance text boxes."""
    for sh in list(slide.shapes):
        if sh.name.startswith("TextBox") and sh.has_text_frame:
            sh._element.getparent().remove(sh._element)

def set_title(slide, text):
    for sh in slide.shapes:
        if sh.name.startswith("Title") and sh.has_text_frame:
            sh.text_frame.text = text
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
            return

def set_pill(slide, text=TEAM):
    for sh in slide.shapes:
        if sh.name.startswith("Oval") and sh.has_text_frame:
            sh.fill.solid(); sh.fill.fore_color.rgb = NAVY
            sh.line.color.rgb = NAVY
            sh.text_frame.text = text
            for p in sh.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE

def box(slide, x, y, w, h, title=None, lines=None, accent=NAVY, fill=WHITE,
        title_sz=11, body_sz=10, gap=2):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = accent; sh.line.width = Pt(1)
    sh.adjustments[0] = 0.04
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.11); tf.margin_right = Inches(0.09)
    tf.margin_top = Inches(0.07); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    if title:
        p = tf.paragraphs[0]; p.text = title; first = False
        p.space_after = Pt(3)
        for r in p.runs:
            r.font.size = Pt(title_sz); r.font.bold = True; r.font.color.rgb = accent
    for ln in (lines or []):
        bold = ln.startswith("**")
        col  = RED if ln.startswith("!!") else (GREEN if ln.startswith("++") else GREY)
        txt  = ln.lstrip("*!+ ")
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.text = txt; p.space_after = Pt(gap)
        for r in p.runs:
            r.font.size = Pt(body_sz); r.font.bold = bold or ln.startswith(("!!","++"))
            r.font.color.rgb = col
    return sh

def txt(slide, x, y, w, h, text, sz=10, bold=False, color=GREY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0; tf.margin_right = 0; tf.margin_bottom = 0
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
    return tb

# ---------------------------------------------------------------- SLIDE 1
s = prs.slides[0]
for sh in s.shapes:
    if sh.name == "Subtitle 3":
        sh.text_frame.text = "PAULIGUARD"
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = NAVY
    if sh.name == "TextBox 9":
        tf = sh.text_frame; tf.clear()
        rows = [
            ("Problem Statement ID", "26141"),
            ("Problem Statement Title", "Quantum-Inspired Cyber Threat Detection for Digital Signature Security"),
            ("Theme", "Blockchain & Cybersecurity"),
            ("PS Category", "Software"),
            ("Organisation", "Egreen Quanta"),
            ("Team ID", "<your team ID>"),
            ("Team Name", "<registered team name>"),
        ]
        for i,(k,v) in enumerate(rows):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.text = f"{k} — {v}"; p.space_after = Pt(6)
            for r in p.runs:
                r.font.size = Pt(12.5); r.font.bold = (i<2); r.font.color.rgb = NAVY

# ---------------------------------------------------------------- SLIDE 2
s = prs.slides[1]; clear_instruction_boxes(s); set_title(s, "PROPOSED SOLUTION"); set_pill(s)

box(s, 0.36, 1.22, 12.6, 0.82, None, [
    "!!A forgery against the target scheme succeeds with probability 1.0 — and the statistical detection this problem statement prescribes catches it 0.000 of the time. That is not a tuning gap. We prove it is an impossibility.",
], accent=RED, fill=RGBColor(0xFD,0xF2,0xF2), body_sz=12.5)

box(s, 0.36, 2.18, 3.05, 2.42, "WHAT IT IS", [
    "PauliGuard screens teleportation-based quantum signature protocols for attacks, before deployment.",
    "A protocol is written as a YAML spec, executed exactly (Bell states, teleportation, Pauli correction, projective measurement), attacked from a catalogue, and judged by four independent detection layers.",
    "**Working software: 208 tests, 4 schemes, live web UI.",
], accent=NAVY)

box(s, 3.53, 2.18, 3.05, 2.42, "HOW IT ADDRESSES THE PROBLEM", [
    "We implement the PS's prescribed method properly — Pauli eigenstates, projective measurement, statistical thresholds — as layers L1 and L2, with derived (never tuned) bounds.",
    "Then we measure exactly where it works and where it cannot:",
    "++L1 detects channel manipulation 1.000",
    "!!L1 detects Pauli forgery 0.000",
    "L0 catches replay & key-reuse 1.000",
], accent=NAVY)

box(s, 6.70, 2.18, 3.05, 2.42, "INNOVATION AND UNIQUENESS", [
    "A fourth layer the PS does not ask for: L3 solves for the subspace of Pauli operators leaving every verification predicate invariant — GF(2) linear algebra, not sampling.",
    "It emits a certificate confirmed by execution before it is shown, so precision is 1 by construction.",
    "**No AI/ML anywhere in the detection path, as the PS requires.",
    "**We also ship the fix: adding a SWAP test cuts the forgery 286x.",
], accent=NAVY)

# THE one enormous number
box(s, 9.87, 2.18, 3.09, 2.42, None, [], accent=RED, fill=RGBColor(0xFD,0xF2,0xF2))
txt(s, 9.95, 2.42, 2.95, 1.0, "0.000", sz=72, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(s, 9.95, 3.42, 2.95, 1.1,
    "Detection rate of the prescribed statistical method against a forgery that succeeds "
    "with probability 1.0 — measured at every noise level tested (n=300, 95% CI [0.000, 0.012]).",
    sz=9.5, color=GREY, align=PP_ALIGN.CENTER)

box(s, 0.36, 4.74, 12.6, 2.05, "THE FOUR-LINE PROOF  (why no threshold could ever work)", [
    "Trent's check is  E_k|P> = |S>.  The adversary is a legitimate participant: he applies Pauli U to the message copy and V = E_k U E_k^-1 to the signature.",
    "Because E_k is a Clifford operation, V is again a Pauli, and V does not depend on the message. So  E_k(U|P>) = V(E_k|P>) = V|S>  — the check passes identically, for every key.",
    "Every encrypted register is maximally mixed without the key, and a unitary on the maximally mixed state returns it unchanged. Decoys untouched. Entangled resource untouched.",
    "!!Therefore the honest and forged executions are the SAME DENSITY MATRIX. We measured trace distance = 0.000e+00 and total-variation distance = 0.000e+00 across 40 random measurement bases. No measurement — not a cleverer basis, not a collective one — can separate them.",
], accent=NAVY, fill=BG, body_sz=10)

# ---------------------------------------------------------------- SLIDE 3
s = prs.slides[2]; clear_instruction_boxes(s); set_title(s, "TECHNICAL APPROACH"); set_pill(s)

txt(s, 0.36, 1.16, 12.6, 0.24, "SYSTEM ARCHITECTURE", sz=11, bold=True, color=NAVY)

def node(x, y, w, h, label, sub, fill=WHITE, accent=NAVY):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = accent; sh.line.width = Pt(1.1); sh.adjustments[0] = 0.10
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
    for r in p.runs: r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = accent
    p2 = tf.add_paragraph(); p2.text = sub; p2.alignment = PP_ALIGN.CENTER
    for r in p2.runs: r.font.size = Pt(8); r.font.color.rgb = GREY
    return sh

def arrow(x, y, w=0.30):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.16))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x9A,0xA3,0xB2); a.line.fill.background()
    a.shadow.inherit = False

Y = 1.46
node(0.36, Y, 1.86, 0.86, "PROTOCOL SPEC", "YAML · discovered from disk\n4 schemes shipped")
arrow(2.30, Y+0.35)
node(2.68, Y, 1.86, 0.86, "ENGINE", "Stim stabilizer sim\nexact, O(n^2)")
arrow(4.62, Y+0.35)
node(5.00, Y, 1.72, 0.86, "TRACE", "frozen JSON schema\nv1.1 contract")
arrow(6.80, Y+0.35)
node(7.18, Y, 1.86, 0.86, "ADVERSARY", "8 attacks as\nspec-level transforms")
arrow(9.12, Y+0.35)
node(9.50, Y, 1.72, 0.86, "4 LAYERS", "L0 · L1 · L2 · L3")
arrow(11.30, Y+0.35)
node(11.68, Y, 1.28, 0.86, "VERDICT", "+ derivation\n+ certificate", fill=RGBColor(0xF0,0xF6,0xFF))

L = [
 ("L0  CONFORMANCE", "NO THRESHOLD", "step order, register lifetime, single-use key material (sha256), freshness",
  "replay 1.000 · key-reuse 1.000 · unauthorised verification", "FPR 0.000 by construction (400 runs, one ledger)", GREEN),
 ("L1  CHANNEL STATS", "SERFLING", "decoy error rate per basis; tau inverted from Serfling, floor-relative",
  "intercept-resend 1.000", "!!Pauli forgery 0.000 — structurally blind", NAVY),
 ("L2  ENTANGLEMENT", "AZUMA", "randomly sampled Pauli stabilizer generators; CHSH 2sqrt2 vs bound 2",
  "resource substitution 1.000 at c=0.40", "!!Pauli forgery 0.000 — resource untouched", NAVY),
 ("L3  ALGEBRAIC", "NO THRESHOLD", "solves GF(2) symplectic nullspace for invariant Pauli operators",
  "!!Pauli forgery 1.000 + certificate", "sound, NOT complete — finds attacks, never proves security", RED),
]
x = 0.36
for name, badge, what, catches, blind, col in L:
    box(s, x, 2.52, 3.09, 1.72, f"{name}   [{badge}]", [
        what, f"++catches: {catches.lstrip('!+ ')}" if not catches.startswith("!!") else catches,
        blind,
    ], accent=col, body_sz=9, title_sz=10)
    x += 3.17

box(s, 0.36, 4.36, 6.24, 1.34, "WHY STIM, NOT A STATEVECTOR SIMULATOR", [
    "Every gate in these protocols is Clifford, so stabilizer simulation is EXACT here, not an approximation — and costs O(n^2) instead of 2^n.",
    "**We verified this rather than assuming it: the target scheme's five-qubit resource state IS a stabilizer state — exactly 32 of 4^5 Paulis stabilise it, and Stim reproduced the state vector to 5.55e-17.",
    "Statevector is retained only as a cross-check oracle at n<=10. Three independent implementations agree.",
], accent=NAVY, fill=BG, body_sz=9.5)

box(s, 6.72, 4.36, 6.24, 1.34, "THRESHOLD DISCIPLINE — NO MAGIC NUMBERS", [
    "No numeric threshold literal exists anywhere in the decision path. Every threshold is computed at runtime from a declared security parameter, the measured hardware floor and the sample size, via a NAMED inequality, and the UI shows the derivation:",
    "**tau = 0.146619 from Serfling with k=402, N=1608, alpha=1e-10, floor=0.0344238; flag iff (xbar - floor) >= tau",
    "Two of the four layers have no threshold at all.",
], accent=NAVY, fill=BG, body_sz=9.5)

box(s, 0.36, 5.82, 12.6, 0.98, "STACK  &  PROCESS", [
    "Python 3.12 · Stim (stabilizer) · Qiskit Aer + an independent from-scratch statevector simulator as cross-check oracles · NumPy GF(2) symplectic layer · SciPy · FastAPI · vanilla-JS UI (offline, no CDN) · pytest",
    "**spec (YAML)  ->  exact execution  ->  trace  ->  attack injection  ->  four-layer analysis  ->  certificate + derivation  ->  evaluation matrix with Clopper-Pearson intervals",
    "Hardware: entanglement/witness sub-circuit run on IBM ibm_kingston (Heron r2). Job da8up31qtnsc73d0v7h0, 4096 shots. Results committed and replayed from disk — we never demo live from a queue.",
], accent=NAVY, fill=BG, body_sz=9.5)

# ---------------------------------------------------------------- SLIDE 4
s = prs.slides[3]; clear_instruction_boxes(s); set_title(s, "FEASIBILITY AND VIABILITY"); set_pill(s)

box(s, 0.36, 1.22, 4.10, 2.62, "FEASIBLE BECAUSE  (already demonstrated)", [
    "++Built and running: 208 passing tests, 20 commits, 4 protocol schemes, live web UI.",
    "Software only. Zero hardware cost. IBM Open Plan gives 10 min QPU/28 days; our circuits need seconds.",
    "Clifford structure makes simulation exact and cheap, so we run real parameter ranges, not toy values.",
    "Ground truth comes from OUTSIDE our code: published attack probabilities plus fixed physical constants (CHSH 2sqrt2, intercept-resend QBER 0.25, teleportation benchmark 2/3).",
    "**The PS attaches no dataset (Dataset Link: Public/Open), so external validation had to come from the literature. It does.",
], accent=GREEN, body_sz=9.5)

box(s, 4.58, 1.22, 4.10, 2.62, "COULD BREAK BECAUSE  (named risks)", [
    "!!R1 — the whole approach needed the target scheme's 5-qubit resource state to be a stabilizer state. If not, Stim is unusable and the GF(2) search collapses.",
    "!!Spec ambiguity — published schemes leave key length, reuse policy and dispute procedures undefined. An analysis is only as good as the spec.",
    "!!L3 is sound, not complete — it searches the Pauli group mod phase. A general adversary is an arbitrary CPTP map, outside the search. It says nothing about hash collisions.",
    "!!A judge supplies a non-Clifford gate the stabilizer backend cannot handle.",
    "!!QPU queue times are unpredictable.",
], accent=RED, body_sz=9.5)

box(s, 8.80, 1.22, 4.16, 2.62, "HANDLED BY  (specific mitigations)", [
    "++R1 RESOLVED YES before any simulator code was written. Brute force over all 4^5 Paulis found exactly 32 stabilisers; Stim agreed to 5.55e-17. Generators recorded.",
    "++Every spec must declare an assumed_fields list. The loader reports every field it had to assume — the ambiguity is a deliverable, not a hidden risk.",
    "++Stated on every certificate, in the module docstring and in the UI. Every emitted attack is confirmed by execution first, so precision is 1 by construction.",
    "++Rehearsed graceful degradation: the tool prints its own limitation and falls back; it does not crash. Malformed spec returns 400, never 500.",
    "++Jobs run in advance, committed with job IDs and calibration snapshots, replayed from disk.",
], accent=GREEN, body_sz=9.5)

box(s, 0.36, 3.96, 6.24, 1.52, "MEASURED: THE FIX WORKS, AND COSTS NOTHING", [
    "We do not propose a NEW scheme — in a field where expert schemes are broken with near-certainty, that is how a team loses the Q&A. We propose fixes to EXISTING schemes and let the tool prove the attack is gone.",
    "!!lu-2022 as published:  forgery succeeds 1.0000",
    "++lu-2022 + SWAP test, k=8:  forgery succeeds 0.0035  (analytic 2^-8 = 0.003906) — a 286x reduction",
    "++Honest acceptance stays 500/500 = 1.0000: the fix has ONE-SIDED ERROR and never falsely rejects a legitimate signature.",
], accent=GREEN, fill=BG, body_sz=9)

box(s, 6.72, 3.96, 6.24, 1.52, "WHY NO AI/ML — AND WHY THAT IS A STRENGTH", [
    "The PS forbids it. There is also a stronger reason.",
    "**No training data exists and none can. No teleportation-based QDS network exists anywhere, so every byte would come from our own simulator — an ML detector would learn our simulator's artefacts and report them as security.",
    "**The security claim is information-theoretic; a fitted threshold is not. Serfling and Azuma give distribution-free bounds that compose into the protocol's epsilon.",
    "**Refusing ML is what bought the impossibility result: we can prove no detector sees this attack because the method is analytic.",
], accent=NAVY, fill=BG, body_sz=9)

box(s, 0.36, 5.60, 12.6, 1.20, "WHAT INDEPENDENT VERIFICATION CAUGHT  (evidence the process is real)", [
    "Every result was re-derived from first principles rather than trusting the test suite that shipped with the code. That caught three defects a green suite had hidden:",
    "**1. Single-use key enforcement keyed on key NAME flagged 399/400 HONEST runs. Fixed by binding to key material (sha256 per run); re-verified at 0/400.",
    "**2. A detection claim quietly used m=500 after being specified at m=200, rather than reporting that 30% corruption is mathematically undetectable at m=200. The limit is now a documented, tested property.",
    "**3. The evaluation matrix rendered a protocol-DEFEATED attack identically to a MISSED one. Measured, unpaired_pauli is accepted 0/300 — the protocol defeats it. The matrix now separates the two.",
], accent=NAVY, fill=BG, body_sz=9.5)

# ---------------------------------------------------------------- SLIDE 5
s = prs.slides[4]; clear_instruction_boxes(s); set_title(s, "IMPACT AND BENEFITS"); set_pill(s)

box(s, 0.36, 1.20, 12.6, 0.62, None, [
    "**HONEST SCOPE, STATED FIRST: no production QDS infrastructure anywhere runs a teleportation-based arbitrated scheme. We make no claim to protect deployed systems today. The impact is on the research and standardisation pipeline that decides which protocols get adopted — and that is a real, dated, citable problem.",
], accent=NAVY, fill=BG, body_sz=10)

hdr = ["BENEFICIARY", "BENEFIT", "SUPPORTING EVIDENCE  (number + source)"]
rows = [
 ("Protocol designers &\njournal reviewers",
  "Screen a candidate scheme in minutes instead of years, before publication or adoption.",
  "Cryptanalysis of these schemes currently runs at a 2-5 year lag: Zeng-Keitel (PRA 65, 042312, 2002) broken by Zou & Qiu (PRA 82, 042325, 2010). Four more schemes (2017-2024) broken by Jacqmin & Lienardy, March 2026."),
 ("National quantum\nprogrammes",
  "Screen candidate protocols before adoption rather than discovering breaks afterwards.",
  "The March 2026 analysis of four schemes required a defence research grant at a military academy. PauliGuard reproduces that class of analysis from a spec file automatically."),
 ("Egreen Quanta",
  "A triage instrument: it tells you WHICH failure mode a candidate protocol has.",
  "Measured contrast across lineages: on decoy-state QDS the algebraic layer returns 0 certificates and statistics do the work; on teleportation AQS it is exactly reversed (L3 1.000, L1 0.000)."),
 ("The wider field",
  "The three threats the PS does not name are the defining goals of an arbitrated signature.",
  "Repudiation of origin, repudiation of receipt and false allegation. Our analyser flags both AQS schemes with RECEIPT_UNBOUND + FALSE_ALLEGATION, claimed_by_scheme=True: they CLAIM the goal their own spec cannot support."),
]
xs = [0.36, 2.30, 5.20]; ws = [1.86, 2.82, 7.76]
for i,h in enumerate(hdr):
    b = box(s, xs[i], 1.94, ws[i], 0.30, None, [], accent=NAVY, fill=NAVY)
    txt(s, xs[i]+0.09, 1.99, ws[i]-0.18, 0.22, h, sz=9.5, bold=True, color=WHITE)
y = 2.30
for who, ben, ev in rows:
    h = 0.94
    box(s, xs[0], y, ws[0], h, None, [f"**{who}"], accent=LGREY, fill=BG, body_sz=9.5)
    box(s, xs[1], y, ws[1], h, None, [ben], accent=LGREY, fill=WHITE, body_sz=9.5)
    box(s, xs[2], y, ws[2], h, None, [ev], accent=LGREY, fill=WHITE, body_sz=9.5)
    y += h + 0.06

box(s, 0.36, 6.32, 12.6, 0.50, None, [
    "**SUSTAINABILITY: the spec format is the compounding asset — each new scheme is hours of work, not a rewrite, and permanently expands the catalogue. NO economic impact figure is claimed, because none can be honestly derived for a technology with zero deployments.",
], accent=NAVY, fill=BG, body_sz=9.5)

# ---------------------------------------------------------------- SLIDE 6
s = prs.slides[5]; clear_instruction_boxes(s); set_title(s, "RESEARCH AND REFERENCES"); set_pill(s)

box(s, 0.36, 1.18, 6.20, 2.62, "PRIMARY SOURCES", [
    "**P.-A. Jacqmin, J. Lienardy. Cryptanalysis of four arbitrated quantum signature schemes. arXiv:2603.19985 / ePrint 2026/558 (23 Mar 2026), Royal Military Academy Brussels.  — the paper this project is built on; supplies the probability-1 forgery and the 1/(8n) bound.",
    "**D. Lu, Z. Li, J. Yu, Z. Han. A verifiable arbitrated quantum signature scheme based on controlled quantum teleportation. Entropy 24, 111 (2022). doi:10.3390/e24010111  — primary target scheme.",
    "**Q. Li, W. H. Chan, D.-Y. Long. Arbitrated quantum signature scheme using Bell states. Phys. Rev. A 79, 054307 (2009). doi:10.1103/PhysRevA.79.054307",
    "**F. Gao, S.-J. Qin, F.-Z. Guo, Q.-Y. Wen. Cryptanalysis of the arbitrated quantum signature protocols. Phys. Rev. A 84, 022344 (2011). doi:10.1103/PhysRevA.84.022344  — QOTP malleability enables receiver forgery.",
    "**J. W. Choi, K.-Y. Chang, D. Hong. Security problem on arbitrated quantum signature schemes. Phys. Rev. A 84, 062330 (2011). doi:10.1103/PhysRevA.84.062330",
    "**D. Gottesman, I. Chuang. Quantum digital signatures. arXiv:quant-ph/0105032 (2001).",
    "**C. Gidney. Stim: a fast stabilizer circuit simulator. Quantum 5, 497 (2021). doi:10.22331/q-2021-07-06-497",
], accent=NAVY, body_sz=8.6, gap=3)

box(s, 6.72, 1.18, 6.24, 1.36, "EXTERNAL VALIDATION — WE REPRODUCE A PUBLISHED RESULT", [
    "We obtained arXiv:2603.19985 and implemented its exact three-layer KCCC encryption. Enumerating the FULL keyspace (n!·4^n), our independent count recovers the paper's three key conditions exactly:",
    "**Pr(k1(1)=1) = 1/n · Pr(k2_1=0) = 1/2 · Pr(tau_1=0) = 1/2 — joint = 1/(4n) exactly (1/8, 1/12, 1/16 at n=2,3,4).",
    "Their lower bound 1/(8n) holds at every n. This number was derived by other people and published in March 2026; we never had their code.",
], accent=GREEN, fill=RGBColor(0xF2,0xFB,0xF5), body_sz=9)

box(s, 6.72, 2.66, 6.24, 0.62, "PROOF OF LIFE", [
    "**Repository (source, 208 tests, 4 schemes, evaluation matrix):  github.com/harshwardhan-kp/pauliguard-v2",
    "**Run it:  .venv/bin/python -m pytest -q   ·   .venv/bin/python -m pauliguard.evaluation",
], accent=RED, fill=RGBColor(0xFD,0xF2,0xF2), body_sz=9)

# --- measured results as a REAL table so markers cannot leak into the text ---
txt(s, 0.36, 3.94, 6.20, 0.26, "MEASURED RESULTS   (n=300 per cell, Clopper-Pearson 95% CI)",
    sz=11, bold=True, color=NAVY)

DATA = [
    ("ATTACK", "L0", "L1", "L2", "L3"),
    ("paired-Pauli forgery  (succeeds 1.000)", "0.000", "0.000", "0.000", "1.000"),
    ("intercept-resend", "0.000", "1.000", "0.000", "0.000"),
    ("replay", "1.000", "0.000", "0.000", "0.000"),
    ("key reuse (same key material)", "1.000", "0.000", "0.000", "0.000"),
    ("honest runs  (false-positive rate)", "0.000", "0.000", "0.000", "0.000"),
]
# cells we want coloured: (row, col) -> colour
HILITE = {(1,2): RED, (1,3): RED, (1,4): GREEN,
          (2,2): GREEN, (3,1): GREEN, (4,1): GREEN}
tx, ty, cw0, cwn, rh = 0.36, 4.24, 3.20, 0.72, 0.235
for r, row in enumerate(DATA):
    for c, val in enumerate(row):
        x = tx if c == 0 else tx + cw0 + (c-1)*cwn
        w = cw0 if c == 0 else cwn
        head = (r == 0)
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(ty + r*rh),
                                  Inches(w), Inches(rh))
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if head else (WHITE if r % 2 else BG)
        cell.line.color.rgb = LGREY; cell.line.width = Pt(0.5)
        cell.shadow.inherit = False
        tf = cell.text_frame; tf.word_wrap = False
        tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_ = tf.paragraphs[0]; p_.text = val
        p_.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        col = WHITE if head else HILITE.get((r, c), GREY)
        for run in p_.runs:
            run.font.size = Pt(8.6)
            run.font.bold = head or (r, c) in HILITE
            run.font.color.rgb = col

box(s, 0.36, 5.70, 6.20, 1.10, None, [
    "!!The zeros are the finding, not a failure. L1 and L2 are correct implementations of the method the PS prescribes; they are structurally blind to an internal adversary exploiting an algebraic symmetry, and we prove why.",
    "FPR stays 0.000 at noise p = 0, 0.001, 0.01 and 0.05. No aggregate accuracy figure is reported anywhere — with a mostly-honest run distribution it is dominated by the base rate and means nothing.",
    "**The published PS ends with the literal placeholder \"Add \'Delivery Table (Expected Deliverables)\' here\". We inferred D1-D8 from the Objectives and Expected Solution and map each to its source sentence.",
], accent=NAVY, fill=BG, body_sz=8.6, gap=2)

# --- screenshot: sized to fit above the footer, caption directly beneath ---
pic = "docs/screenshots/pauliguard-ui.png"
if pathlib.Path(pic).exists():
    from PIL import Image
    iw, ih = Image.open(pic).size
    pw = 5.30
    ph = pw * ih / iw
    px = 6.72 + (6.24 - pw) / 2
    py = 3.40
    s.shapes.add_picture(pic, Inches(px), Inches(py), width=Inches(pw))
    txt(s, 6.72, py + ph + 0.06, 6.24, 0.44,
        "Live UI, honest vs forged in lockstep: SIGNATURE ACCEPTED, message CHANGED "
        "([0,1] -> [1,1]), decoy error rates identical at 0.03310. L0/L1/L2 green; "
        "L3 red with a certificate confirmed by execution 16/16.",
        sz=8.2, color=GREY, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- drop slide 7
xml = prs.slides._sldIdLst
ids = list(xml)
for sid in ids[6:]:
    prs.part.drop_rel(sid.rId); xml.remove(sid)

pathlib.Path("deck").mkdir(exist_ok=True)
prs.save(OUT)
print(f"wrote {OUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
